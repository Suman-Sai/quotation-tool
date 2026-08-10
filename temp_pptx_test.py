from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import warnings
warnings.simplefilter('always')
prs = Presentation()
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
background = slide.background
fill = background.fill
fill.gradient()
fill.gradient_stops[0].color.rgb = RGBColor(135, 206, 235)
fill.gradient_stops[1].color.rgb = RGBColor(30, 144, 255)
title_shape = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(1))
title_frame = title_shape.text_frame
title_frame.text = 'Weather Forecast Dashboard'
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle_shape = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(1))
subtitle_frame = subtitle_shape.text_frame
subtitle_frame.text = 'Powered by Power BI & WeatherAPI'
subtitle_frame.paragraphs[0].font.size = Pt(24)
subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
prs.save('Weather_Forecast_Theme.pptx')
print('done')
